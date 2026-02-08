import argparse
import json
import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt

def create_presentation(data):
    prs = Presentation()
    
    # Extract general info
    output_file = data.get("output_file", "presentation.pptx")
    slides_data = data.get("slides", [])
    
    # Layout indices for standard template
    LAYOUT_TITLE = 0
    LAYOUT_TITLE_CONTENT = 1
    LAYOUT_SECTION = 2
    
    print(f"Generating presentation with {len(slides_data)} slides...")

    for i, slide_info in enumerate(slides_data):
        layout_name = slide_info.get("layout", "Title and Content")
        title_text = slide_info.get("title", "")
        
        # Determine layout
        if layout_name.lower() == "title slide":
            slide_layout = prs.slide_layouts[LAYOUT_TITLE]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set Title
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            
            # Set Subtitle
            subtitle_text = slide_info.get("subtitle", "")
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = subtitle_text
                
        elif layout_name.lower() == "section header":
            slide_layout = prs.slide_layouts[LAYOUT_SECTION]
            slide = prs.slides.add_slide(slide_layout)
             # Set Title
            if slide.shapes.title:
                slide.shapes.title.text = title_text
                
        else: # Default to Title and Content
            slide_layout = prs.slide_layouts[LAYOUT_TITLE_CONTENT]
            slide = prs.slides.add_slide(slide_layout)
            
            # Set Title
            if slide.shapes.title:
                slide.shapes.title.text = title_text
            
            # Set Content
            content_data = slide_info.get("content", "")
            
            # Find the content placeholder (usually index 1)
            # A more robust way is to find a placeholder that contains text but isn't title
            body_shape = None
            for shape in slide.placeholders:
                if shape.shape_id != slide.shapes.title.shape_id:
                     if hasattr(shape, "text"):
                        body_shape = shape
                        break
            
            if body_shape:
                tf = body_shape.text_frame
                tf.clear() # Clear existing placeholder text
                
                if isinstance(content_data, list):
                    for idx, item in enumerate(content_data):
                        p = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
                        p.text = str(item)
                        p.level = 0
                else:
                    tf.text = str(content_data)

    # Save
    try:
        prs.save(output_file)
        print(f"Presentation saved to: {os.path.abspath(output_file)}")
    except Exception as e:
        print(f"Error saving presentation: {e}", file=sys.stderr)
        sys.exit(1)

def main():
    parser = argparse.ArgumentParser(description="Generate PowerPoint from JSON input.")
    group = parser.add_mutually_exclusive_group(required=True)
    group.add_argument("--json", type=str, help="JSON string input")
    group.add_argument("--file", type=str, help="Path to JSON file")
    
    args = parser.parse_args()
    
    data = None
    try:
        if args.file:
            with open(args.file, 'r', encoding='utf-8') as f:
                data = json.load(f)
        else:
            data = json.loads(args.json)
            
    except json.JSONDecodeError as e:
        print(f"Error decoding JSON: {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"Error reading input: {e}", file=sys.stderr)
        sys.exit(1)
        
    create_presentation(data)

if __name__ == "__main__":
    main()
